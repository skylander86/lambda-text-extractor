# -*- coding: utf-8 -*-

import codecs
from datetime import datetime, time
import json
import logging
import os
import subprocess
from tempfile import NamedTemporaryFile
from time import sleep

import boto3
import pptx
from PyPDF2 import PdfFileReader
import xlrd

from utils import download_file, upload_file, get_file_content, delete_objects, head_object

LAMBDA_TASK_ROOT = os.environ.get('LAMBDA_TASK_ROOT', os.path.dirname(os.path.abspath(__file__)))
BIN_DIR = os.path.join(LAMBDA_TASK_ROOT, 'bin')
LIB_DIR = os.path.join(LAMBDA_TASK_ROOT, 'lib')

logging.basicConfig(format=u'%(asctime)-15s [%(name)s] %(levelname)s: %(message)s', level=logging.INFO)
logger = logging.getLogger()

lambda_client = boto3.client('lambda')
s3_client = boto3.client('s3')


def extract(event, context):
    doc_uri = event['doc_uri']
    text_uri = event['text_uri']
    text_encoding = event.get('text_encoding', 'utf-8')

    try: doc_path = download_file(doc_uri)
    except Exception as e: return dict(success=False, reason=u'Exception while downloading from <{}>: {}'.format(doc_uri, e))

    _, ext = os.path.splitext(doc_path)  # get format from extension
    parse_func = PARSE_FUNCS.get(ext.lower())

    if parse_func is None:
        return dict(success=False, reason=u'Unknown file type <{}>'.format(doc_uri))

    if ext == '.pdf' and event.get('force_ocr'):
        parse_func = pdf_to_text_with_ocr_single_page if 'page' in event else pdf_to_text_with_ocr
    #end if

    o = parse_func(doc_path, event, context)
    # try: o = parse_func(doc_path, event, context)
    # except Exception as e: return dict(success=False, reason=u'Exception while parsing <{}> using <{}>: {}'.format(doc_uri, parse_func.__name__, e))

    if not o['success']: return o

    text = o['text']

    with NamedTemporaryFile(prefix='text-extractor.', suffix='.txt', delete=False) as f:
        text_path = f.name
        if isinstance(text, unicode):
            f.write(text.encode(text_encoding))
        else:
            logger.warning(u'Text extracted from <{}> is not in unicode!'.format(doc_uri))
            f.write(text)
        #end if
    #end with

    try: upload_file(text_uri, text_path)
    except Exception as e: return dict(success=False, reason=u'Exception while uploading to <{}>: {}'.format(text_uri, e))

    if not event.get('text', False): del o['text']
    o['doc_uri'] = doc_uri
    o['text_uri'] = text_uri
    o['size'] = len(text)
    o.setdefault('method', parse_func.__name__)

    return o
#end def


def pdf_to_text(doc_path, event, context):
    cmdline = [os.path.join(BIN_DIR, 'pdftotext'), '-nopgbrk', doc_path]
    try:
        subprocess.check_call(cmdline, shell=False, stderr=subprocess.STDOUT)
        text_path = os.path.splitext(doc_path)[0] + '.txt'
    except subprocess.CalledProcessError as e: return dict(success=False, reason=u'Exception while executing {}: {} (output={})'.format(cmdline, e, e.output))

    with codecs.open(text_path, 'r', 'utf-8', errors='ignore') as f:
        text = f.read()
    text = text.strip()

    if len(text) < 512:
        return pdf_to_text_with_ocr(doc_path, event, context)

    return dict(success=True, text=text)
#end def


def pdf_to_text_with_ocr(doc_path, event, context):
    with open(doc_path, 'rb') as f:
        reader = PdfFileReader(f)
        num_pages = reader.getNumPages()
    #end with

    doc_uri = event['doc_uri']
    text_uri = event['text_uri']

    page_text_uris = {}
    for page in xrange(1, num_pages + 1):
        page_text_uri = text_uri + '-{}'.format(page)
        lambda_client.invoke(
            FunctionName='text-extractor_text_extractor',
            InvocationType='Event',
            LogType='None',
            Payload=json.dumps(
                dict(doc_uri=doc_uri, text_uri=page_text_uri, page=page, force_ocr=True)
            )
        )
        page_text_uris[page_text_uri] = page
    #end for

    processing = page_text_uris.copy()
    page_contents = {}
    while len(processing) > 0 and context.get_remaining_time_in_millis() > 5000:  # leave 10 seconds to concatenate text and return it
        sleep(1.0)

        for uri, page in processing.items():
            if head_object(uri):
                page_contents[page] = get_file_content(uri)
                del processing[uri]
            #end if
        #end for
    #end while

    delete_objects(page_text_uris.keys())

    contents = []
    for page in xrange(1, num_pages + 1):
        content = page_contents.get(page)
        if content: contents.append(content)
        else: logger.warn(u'Unable to OCR page {} of <{}>.'.format(page, doc_uri))
    #end for

    return dict(success=True, text=u'\n\n'.join(contents), method='pdf_to_text_with_ocr')
#end def


def pdf_to_text_with_ocr_single_page(doc_path, event, context):
    pageno = event['page']
    with NamedTemporaryFile(prefix='text-extractor.', suffix='.png', delete=False) as f:
        image_page_path = f.name

    cmdline = [os.path.join(BIN_DIR, 'gs-920-linux_x86_64'), '-sDEVICE=png16m', '-dFirstPage={}'.format(pageno), '-dLastPage={}'.format(pageno), '-dINTERPOLATE', '-r300', '-o', image_page_path, '-dNOPAUSE', '-dSAFER', '-c', '67108864', 'setvmthreshold', '-dGraphicsAlphaBits=4', '-dTextAlphaBits=4', '-f', doc_path]
    try: subprocess.check_output(cmdline, shell=False, stderr=subprocess.STDOUT)
    except subprocess.CalledProcessError as e: return dict(success=False, reason=u'Exception while executing {}: {} (output={})'.format(cmdline, e, e.output))

    return image_to_text(image_page_path, event, context)
#end def


def image_to_text(doc_path, event, context):
    _, ext = os.path.splitext(doc_path)
    if ext not in ['.png', '.tiff']:  # convert to B&W tiff
        cmdline = [os.path.join(BIN_DIR, 'magick'), '-depth', '8', '-density', '300', doc_path, '-threshold', '50%', '-channel', 'r', '-separate', doc_path + '.tiff']
        try: subprocess.check_output(cmdline, stderr=subprocess.STDOUT, shell=False, env=dict(LD_LIBRARY_PATH=os.path.join(LIB_DIR, 'tesseract')))
        except subprocess.CalledProcessError as e: return dict(success=False, reason=u'Exception while executing {}: {} (output={})'.format(cmdline, e, e.output))

        # Using PIL takes up more memory but 10% faster.
        # im = Image.open(doc_path)
        # doc_path += '.tiff'
        # im.save(doc_path, dpi=(300, 300))
    #end if

    cmdline = [os.path.join(BIN_DIR, 'tesseract'), doc_path, doc_path, '-l', 'eng', '-psm', '6', '--tessdata-dir', os.path.join(LIB_DIR, 'tesseract')]
    try: subprocess.check_output(cmdline, stderr=subprocess.STDOUT, shell=False, env=dict(LD_LIBRARY_PATH=os.path.join(LIB_DIR, 'tesseract')))
    except subprocess.CalledProcessError as e: return dict(success=False, reason=u'Exception while executing {}: {} (output={})'.format(cmdline, e, e.output))
    with codecs.open(doc_path + '.txt', 'r', 'utf-8') as f:
        text = f.read().strip()

    return dict(success=True, text=text)
#end def


def doc_to_text(doc_path, event, context):
    cmdline = [os.path.join(BIN_DIR, 'antiword'), '-t', '-w', '0', '-m', 'UTF-8', doc_path]
    try:
        text = subprocess.check_output(cmdline, shell=False, stderr=subprocess.STDOUT, env=dict(ANTIWORDHOME=os.path.join(LIB_DIR, 'antiword')))
        text = text.decode('utf-8', errors='ignore')
    except subprocess.CalledProcessError as e: return dict(success=False, reason=u'Exception while executing {}: {} (output={})'.format(cmdline, e, e.output))

    return dict(success=True, text=text.strip())
#end def


def rtf_to_text(doc_path, event, context):
    cmdline = [os.path.join(BIN_DIR, 'unrtf'), '-P', os.path.join(LIB_DIR, 'unrtf'), '--text', doc_path]
    try:
        text = subprocess.check_output(cmdline, shell=False, stderr=subprocess.STDOUT)
        text = text.decode('utf-8', errors='ignore')

        new_lines = []
        in_header = True
        for line in text.split(u'\n'):
            if in_header and line.startswith('###'): continue
            else:
                new_lines.append(line)
                in_header = False
            #end if
        #end for
        text = u'\n'.join(new_lines)
    except subprocess.CalledProcessError as e: return dict(success=False, reason=u'Exception while executing {}: {} (output={})'.format(cmdline, e, e.output))

    return dict(success=True, text=text.strip())
#end def


def xls_to_text(doc_path, event, context):
    book = xlrd.open_workbook(doc_path)
    lines = []
    for sheet in book.sheets():
        lines.append(sheet.name)
        lines.append(u'-------------------------------------------')
        for row in sheet.get_rows():
            row_values = []
            for cell in row:
                if cell.ctype == xlrd.XL_CELL_DATE:
                    d = datetime(*xlrd.xldate_as_tuple(cell.value, book.datemode))
                    row_values.append(d.date() if d.time() == time(0, 0) else d)
                elif cell.ctype == xlrd.XL_CELL_BOOLEAN: row_values.append(bool(cell.value))
                else: row_values.append(cell.value)
            #end for
            lines.append(u'\t|\t'.join(map(lambda s: unicode(s).strip(), row_values)))
    #end for

    return dict(success=True, text=u'\n'.join(lines))
#end def


def pptx_to_text(doc_path, event, context):
    prs = pptx.Presentation(doc_path)

    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            text_runs += [run.text for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]
        #end for
    #end for

    return dict(success=True, text=u'\n\n'.join(text_runs))
#end def


PARSE_FUNCS = {
    '.pdf': pdf_to_text,
    '.png': image_to_text,
    '.tiff': image_to_text,
    '.jpg': image_to_text,
    '.jpeg': image_to_text,
    '.rtf': rtf_to_text,
    '.doc': doc_to_text,
    '.xls': xls_to_text,
    '.xlsx': xls_to_text,
    '.pptx': pptx_to_text,
}

if __name__ == '__main__':
    print pdf_to_text_with_ocr('../../Downloads/image_pdf.pdf', None, None)['text']
#end if
