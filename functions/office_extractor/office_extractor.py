# -*- coding: utf-8 -*-

import codecs
from datetime import datetime, time
import logging
import os
import subprocess
from tempfile import NamedTemporaryFile

from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P

import lxml.html
from lxml import etree

import pptx
import xlrd

from utils import download_file, upload_file

LAMBDA_TASK_ROOT = os.environ.get('LAMBDA_TASK_ROOT', os.path.dirname(os.path.abspath(__file__)))
BIN_DIR = os.path.join(LAMBDA_TASK_ROOT, 'bin')
LIB_DIR = os.path.join(LAMBDA_TASK_ROOT, 'lib')

logging.basicConfig(format=u'%(asctime)-15s [%(name)s] %(levelname)s: %(message)s', level=logging.INFO)
logger = logging.getLogger()


def extract(event, context):
    doc_uri = event['doc_uri']
    text_uri = event['text_uri']
    text_encoding = event.get('text_encoding', 'utf-8')

    try: doc_path = download_file(doc_uri)
    except Exception as e: return dict(success=False, reason=u'Exception while downloading from <{}>: {}'.format(doc_uri, e))

    _, ext = os.path.splitext(doc_path)  # get format from extension
    parse_func = PARSE_FUNCS.get(ext.lower())

    if parse_func is None:
        return dict(success=False, reason=u'Unknown file type <{}>'.format(doc_uri))

    o = parse_func(doc_path, event, context)
    # try: o = parse_func(doc_path, event, context)
    # except Exception as e: return dict(success=False, reason=u'Exception while parsing <{}> using <{}>: {}'.format(doc_uri, parse_func.__name__, e))

    if not o['success']: return o

    text = o['text']

    with NamedTemporaryFile(prefix='office-extractor.', suffix='.txt', delete=False) as f:
        text_path = f.name
        if isinstance(text, unicode):
            f.write(text.encode(text_encoding))
        else:
            logger.warning(u'Text extracted from <{}> is not in unicode!'.format(doc_uri))
            f.write(text)
        #end if
    #end with

    try: upload_file(text_uri, text_path)
    except Exception as e: return dict(success=False, reason=u'Exception while uploading to <{}>: {}'.format(text_uri, e))

    if not event.get('text', False): del o['text']
    o['doc_uri'] = doc_uri
    o['text_uri'] = text_uri
    o['size'] = len(text)
    o.setdefault('method', parse_func.__name__)

    return o
#end def


def doc_to_text(doc_path, event, context):
    cmdline = [os.path.join(BIN_DIR, 'antiword'), '-t', '-w', '0', '-m', 'UTF-8', doc_path]
    try:
        text = subprocess.check_output(cmdline, shell=False, stderr=subprocess.STDOUT, env=dict(ANTIWORDHOME=os.path.join(LIB_DIR, 'antiword')))
        text = text.decode('utf-8', errors='ignore')
        text = text.strip()
    except subprocess.CalledProcessError as e: return dict(success=False, reason=u'Exception while executing {}: {} (output={})'.format(cmdline, e, e.output))

    return dict(success=True, text=text)
#end def


def docx_to_text(doc_path, event, context):
    doc = Document(doc_path)
    doc_body = doc.element.body
    blocks = []
    for child in doc_body.iterchildren():
        if isinstance(child, CT_P): blocks.append(Paragraph(child, doc_body).text)
        elif isinstance(child, CT_Tbl): blocks += [cell.text for row in Table(child, doc_body).rows for cell in row.cells]
    #end for

    text = u'\n\n'.join(blocks).strip()

    return dict(success=True, text=text)
#end def


def rtf_to_text(doc_path, event, context):
    cmdline = [os.path.join(BIN_DIR, 'unrtf'), '-P', os.path.join(LIB_DIR, 'unrtf'), '--text', doc_path]
    try:
        text = subprocess.check_output(cmdline, shell=False, stderr=subprocess.STDOUT)
        text = text.decode('utf-8', errors='ignore')

        new_lines = []
        in_header = True
        for line in text.split(u'\n'):
            if in_header and line.startswith('###'): continue
            else:
                new_lines.append(line)
                in_header = False
            #end if
        #end for
        text = u'\n'.join(new_lines)
    except subprocess.CalledProcessError as e: return dict(success=False, reason=u'Exception while executing {}: {} (output={})'.format(cmdline, e, e.output))

    return dict(success=True, text=text.strip())
#end def


def xls_to_text(doc_path, event, context):
    book = xlrd.open_workbook(doc_path)
    lines = []
    for sheet in book.sheets():
        lines.append(sheet.name)
        lines.append(u'-------------------------------------------')
        for row in sheet.get_rows():
            row_values = []
            for cell in row:
                if cell.ctype == xlrd.XL_CELL_DATE:
                    d = datetime(*xlrd.xldate_as_tuple(cell.value, book.datemode))
                    row_values.append(d.date() if d.time() == time(0, 0) else d)
                elif cell.ctype == xlrd.XL_CELL_BOOLEAN: row_values.append(bool(cell.value))
                else: row_values.append(cell.value)
            #end for
            lines.append(u'\t|\t'.join(map(lambda s: unicode(s).strip(), row_values)))
    #end for

    return dict(success=True, text=u'\n'.join(lines))
#end def


def pptx_to_text(doc_path, event, context):
    prs = pptx.Presentation(doc_path)

    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            text_runs += [run.text for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]
        #end for
    #end for

    return dict(success=True, text=u'\n\n'.join(text_runs))
#end def


def html_to_text(doc_path, event, context):
    document = lxml.html.parse(doc_path)
    text = u'\n'.join(etree.XPath("//text()")(document))

    return dict(success=True, text=text)
#end def


def text_to_text(doc_path, event, context):
    with codecs.open(doc_path, 'r', 'utf-8', errors='ignore') as f:
        text = f.read()
    return dict(success=True, text=text)
#end def


def csv_to_text(doc_path, event, context):
    with codecs.open(doc_path, 'r', 'utf-8', errors='ignore') as f:
        text = f.read()
        text = text.replace(u',', u', ')
    #end with

    return dict(success=True, text=text)
#end def


PARSE_FUNCS = {
    '.doc': doc_to_text,
    '.docx': docx_to_text,
    '.rtf': rtf_to_text,
    '.xls': xls_to_text,
    '.xlsx': xls_to_text,
    '.pptx': pptx_to_text,
    '.html': html_to_text,
    '.htm': html_to_text,
    '.txt': text_to_text,
    '.csv': csv_to_text,
}
